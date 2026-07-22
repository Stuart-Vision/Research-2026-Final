"""
Final Defense presentation generator
====================================
S. Shanthosh | 20APSE4882

Builds the final-defence deck on top of the departmental template
("FINAL Defense PPT Template_18-19(1).pptx") so that the theme, colours
and branding are inherited exactly.

Design constants below were extracted directly from that template:
    canvas   13.33 x 7.5 in (16:9)
    brand    #203864
    titles   Times New Roman 44pt
    subtitle Roboto Medium
    body     Roboto Light
    footers  Times New Roman 14pt
    logo     top-right (10.88, 0.0) 2.45 x 0.78 in

Slide order follows the "Research Project Final Defence - Guidelines":
    Introduction, Literature Review, Research Objectives, Methodology,
    Results, Discussion, Future Works, Publications, References.
"""

import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ----------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(HERE, "FINAL Defense PPT Template_18-19(1).pptx")
LOGO = os.path.join(HERE, "logo.jpg")
CHARTS = os.path.join(ROOT, "results")
OUTPUT = os.path.join(HERE, "20APSE4882.pptx")

# ----------------------------------------------------------------------
# Design constants taken from the template
# ----------------------------------------------------------------------
BRAND = RGBColor(0x20, 0x38, 0x64)      # #203864 departmental navy
BODY_CLR = RGBColor(0x1A, 0x1A, 0x1A)   # near-black body text
MUTED = RGBColor(0x55, 0x5E, 0x70)      # secondary text
ACCENT = RGBColor(0xC0, 0x50, 0x2A)     # highlight for key numbers
GOOD = RGBColor(0x1E, 0x6B, 0x3A)       # positive result
RULE = RGBColor(0xC8, 0xD0, 0xDF)       # hairline rules
BAND = RGBColor(0xEE, 0xF2, 0xF8)       # light callout band

F_TITLE = "Times New Roman"
F_SUB = "Roboto Medium"
F_BODY = "Roboto Light"
F_HEAD = "Roboto Black"

DEPT = "DEPARTMENT OF COMPUTING AND INFORMATION SYSTEMS"
# NOTE: taken verbatim from the departmental template. Change this one
# string if your module code differs.
MODULE = "IS 42853 - RESEARCH PROJECT IN IS"

SW, SH = 13.333, 7.5   # slide size in inches

_page = 0              # running slide number stamped bottom-centre


# ----------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------
def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    """Return a fixed-size, word-wrapping text box (no autofit surprises)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, text, *, size=18, font=F_BODY, color=BODY_CLR, bold=False,
          italic=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0,
          line=1.15, first=False):
    """Append (or fill) a paragraph and return it."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def rich(tf, parts, *, size=18, align=PP_ALIGN.LEFT, space_after=6,
         space_before=0, line=1.15, first=False):
    """Append a paragraph built from (text, font, colour, bold) tuples."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    for text, font, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = font
        r.font.bold = bold
        r.font.color.rgb = color
    return p


def hanging(p, left=0.30):
    """Give a paragraph a hanging indent so wrapped lines clear the bullet."""
    emu = int(Inches(left))
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(emu))
    pPr.set("indent", str(-emu))
    return p


def bullet(tf, text, *, size=18, color=BODY_CLR, bold=False, indent=0,
           space_after=9, marker="▪", hang=None):
    """Template-style square bullet with a hanging indent.

    The indent is scaled to the font size so that wrapped lines line up
    with the first line's text rather than with the bullet glyph.
    """
    p = tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = 1.12
    p.level = indent
    hanging(p, hang if hang is not None else size * 0.021)
    m = p.add_run()
    m.text = marker + "   "
    m.font.size = Pt(size)
    m.font.name = F_BODY
    m.font.color.rgb = BRAND
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = F_BODY
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill=None, line=None, width=1.0):
    """Plain rectangle used for bands, rules and callouts."""
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(width)
    sh.text_frame.word_wrap = True
    return sh


def picture(slide, name, x, y, w=None, h=None):
    """Place a chart from results/ if it exists."""
    path = os.path.join(CHARTS, name)
    if not os.path.exists(path):
        print("   ! missing chart:", name)
        return None
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


def table(slide, x, y, w, h, data, *, col_w=None, head_size=12, body_size=12,
          head_fill=BRAND, zebra=True, row_h=None):
    """Compact table styled to the departmental palette."""
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False

    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)

    for ri, row in enumerate(data):
        tbl.rows[ri].height = Inches(row_h or (0.32 if ri else 0.36))
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = head_fill
            elif zebra and ri % 2 == 0:
                cell.fill.fore_color.rgb = BAND
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            # A leading "*" marks the row that should be emphasised.
            emph = isinstance(val, str) and val.startswith("*")
            if emph:
                val = val[1:]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(head_size if ri == 0 else body_size)
            r.font.name = F_SUB if ri == 0 else F_BODY
            r.font.bold = bool(ri == 0 or emph)
            r.font.color.rgb = (RGBColor(0xFF, 0xFF, 0xFF) if ri == 0
                                else (BRAND if emph else BODY_CLR))
    return tbl


# ----------------------------------------------------------------------
# Slide chrome (logo, footers, page number) — matches the template
# ----------------------------------------------------------------------
def chrome(slide, numbered=True):
    global _page
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(10.88), Inches(0.0),
                                 width=Inches(2.45), height=Inches(0.78))

    tf = textbox(slide, 0.22, 7.10, 6.02, 0.32)
    write(tf, DEPT, size=14, font=F_TITLE, color=BRAND, first=True,
          space_after=0)

    tf = textbox(slide, 9.72, 7.13, 3.44, 0.28)
    write(tf, MODULE, size=14, font=F_TITLE, color=BRAND,
          align=PP_ALIGN.RIGHT, first=True, space_after=0)

    if numbered:
        _page += 1
        tf = textbox(slide, 6.39, 7.13, 0.87, 0.28)
        write(tf, str(_page), size=14, font=F_TITLE, color=BRAND,
              align=PP_ALIGN.CENTER, first=True, space_after=0)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # 6 = blank


def heading(slide, title, subtitle=None, *, size=34):
    """Template heading block: Times New Roman title + Roboto sub-title."""
    tf = textbox(slide, 0.35, 0.20, 10.3, 0.72, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, title.upper(), size=size, font=F_TITLE, color=BRAND,
          first=True, space_after=0)
    rect(slide, 0.35, 0.95, 1.5, 0.045, fill=BRAND)
    if subtitle:
        tf = textbox(slide, 0.35, 1.05, 12.4, 0.44)
        write(tf, subtitle, size=19, font=F_SUB, color=MUTED, first=True,
              space_after=0)
        return 1.62
    return 1.30


def content_slide(prs, title, subtitle=None, **kw):
    s = new_slide(prs)
    chrome(s)
    top = heading(s, title, subtitle, **kw)
    return s, top


# ======================================================================
#  BUILD
# ======================================================================
def build():
    global _page
    prs = Presentation(TEMPLATE)

    # Start from the template's theme but drop its placeholder slides.
    for sld in list(prs.slides._sldIdLst):
        rid = sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships}id")
        prs.part.drop_rel(rid)
        prs.slides._sldIdLst.remove(sld)

    # ------------------------------------------------------------------
    # 1. TITLE
    # ------------------------------------------------------------------
    s = new_slide(prs)
    chrome(s, numbered=False)
    # title slide carries the larger logo, as in the template
    for sh in list(s.shapes):
        if sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(8.22), Inches(0.76),
                             width=Inches(4.24), height=Inches(1.20))

    tf = textbox(s, 0.90, 2.55, 11.5, 1.85, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, "An Empirical Evaluation of Deep Learning Models for "
              "Enhanced and Scalable Bug Detection and Classification "
              "in Large-Scale Software Systems",
          size=32, font=F_TITLE, color=BRAND, align=PP_ALIGN.CENTER,
          first=True, space_after=0, line=1.12)

    rect(s, 5.17, 4.62, 3.0, 0.04, fill=BRAND)

    tf = textbox(s, 2.20, 4.86, 8.9, 0.95)
    write(tf, "S. Shanthosh", size=24, font=F_SUB, color=BRAND,
          align=PP_ALIGN.CENTER, first=True, space_after=3)
    write(tf, "researchp180@gmail.com   |   20APSE4882", size=17,
          font=F_BODY, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)

    tf = textbox(s, 1.01, 6.12, 5.6, 0.80)
    write(tf, "Supervisor:", size=15, font=F_SUB, color=BRAND, first=True,
          space_after=2)
    write(tf, "Mrs. W.M.L.S. Abeythunga", size=15, font=F_BODY,
          color=BODY_CLR, space_after=1)
    write(tf, "Department of Computing and Information Systems",
          size=13, font=F_BODY, color=MUTED, space_after=0)

    # ------------------------------------------------------------------
    # 2. CONTENTS
    # ------------------------------------------------------------------
    s = new_slide(prs)
    chrome(s)
    tf = textbox(s, 0.35, 0.24, 6.0, 0.84, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, "CONTENTS", size=44, font=F_TITLE, color=BRAND, first=True,
          space_after=0)
    rect(s, 0.35, 1.16, 1.5, 0.045, fill=BRAND)

    items = [
        ("01", "Introduction"),
        ("02", "Literature Review"),
        ("03", "Research Objectives"),
        ("04", "Methodology"),
        ("05", "Results"),
        ("06", "Discussion"),
        ("07", "Future Works"),
        ("08", "Publications & References"),
    ]
    y = 1.62
    for i, (num, label) in enumerate(items):
        col = i // 4
        row = i % 4
        x = 1.30 + col * 5.90
        yy = y + row * 1.16
        rect(s, x, yy, 0.62, 0.62, fill=BRAND)
        tf = textbox(s, x, yy + 0.13, 0.62, 0.40)
        write(tf, num, size=17, font=F_SUB, color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER, first=True, space_after=0)
        tf = textbox(s, x + 0.88, yy + 0.12, 4.6, 0.45,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tf, label, size=21, font=F_HEAD, color=BRAND, first=True,
              space_after=0)

    # ------------------------------------------------------------------
    # 3. INTRODUCTION — background & motivation
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Introduction",
                           "Background and motivation")
    tf = textbox(s, 0.60, top, 7.35, 4.9)
    for t in [
        "Modern systems run to millions of lines of code, built by "
        "distributed teams under continuous delivery.",
        "Poor software quality costs the US economy over $2 trillion a "
        "year, much of it from defects that escaped detection "
        "(CISQ, 2022).",
        "A defect found in production costs roughly 10x more to fix than "
        "one found at design time (Boehm, 1981).",
        "Manual review does not scale; static analysis produces false "
        "positives until developers stop reading the warnings.",
        "Machine learning can learn defect patterns from historical "
        "software metrics and flag fault-prone modules before release.",
    ]:
        bullet(tf, t, size=17, space_after=13)

    rect(s, 8.35, top + 0.10, 4.45, 4.30, fill=BAND)
    rect(s, 8.35, top + 0.10, 0.075, 4.30, fill=BRAND)
    tf = textbox(s, 8.68, top + 0.38, 3.85, 3.80)
    write(tf, "WHY IT MATTERS", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=14)
    for big, small in [("$2 trillion", "annual US cost of poor software quality"),
                       ("10x", "cost multiplier for a defect found after release")]:
        write(tf, big, size=30, font=F_SUB, color=BRAND, space_after=3,
              line=1.0)
        write(tf, small, size=13, font=F_BODY, color=MUTED, space_after=22,
              line=1.2)

    # ------------------------------------------------------------------
    # 4. RESEARCH PROBLEM
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Research Problem",
                           "Severe class imbalance defeats standard classifiers")

    rect(s, 0.60, top, 7.10, 1.28, fill=BAND)
    rect(s, 0.60, top, 0.075, 1.28, fill=ACCENT)
    tf = textbox(s, 0.92, top + 0.20, 6.6, 0.95)
    write(tf, "Defective modules are 5-20% of any codebase. A model that "
              "maximises accuracy learns to predict “clean” almost "
              "always — scoring well while finding almost no defects.",
          size=16, font=F_BODY, color=BODY_CLR, first=True, space_after=0,
          line=1.22)

    tf = textbox(s, 0.60, top + 1.52, 7.10, 2.9)
    write(tf, "THE ACCURACY PARADOX ON JM1", size=12, font=F_SUB,
          color=MUTED, first=True, space_after=10)
    for t, c in [
        ("Predict every module clean  →  80.7% accuracy, 0 defects found",
         ACCENT),
        ("Random Forest, no correction  →  81.1% accuracy, 85 of 421 found",
         BODY_CLR),
    ]:
        bullet(tf, t, size=16, color=c, space_after=10)
    write(tf, "0.4 percentage points of accuracy separate a useless model "
              "from a partly useful one. Accuracy cannot rank these models.",
          size=15, font=F_BODY, color=MUTED, italic=True, space_before=8,
          space_after=0, line=1.2)

    picture(s, "01_class_distribution.png", 8.20, top + 0.10, w=4.65)

    # ------------------------------------------------------------------
    # 5. LITERATURE REVIEW
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Literature Review",
                           "IEEE Xplore, ACM DL, Springer and ScienceDirect; "
                           "seminal works plus 2016–2024")
    table(s, 0.60, top + 0.05, 12.15, 4.3,
          [["Study", "Approach", "Contribution", "Limitation"],
           ["Halstead (1977);\nMcCabe (1976)", "Complexity metrics",
            "Established the metric families still used today",
            "Static snapshot only"],
           ["Menzies et al. (2010)", "Traditional ML on\nstatic features",
            "Simple models on good features are hard to beat",
            "Ceiling set by static features"],
           ["Li (2017)", "CNN on code and\nmetric vectors",
            "Improved F1 over traditional ML on PROMISE",
            "Local patterns only"],
           ["Hochreiter &\nSchmidhuber (1997)", "LSTM",
            "Models long-range sequential dependency",
            "No local feature detection"],
           ["Lin et al. (2017)", "Focal Loss",
            "Reweights training toward hard examples",
            "Validated in vision, not in defect prediction"],
           ["Chawla (2002);\nHan (2005)", "SMOTE /\nBorderlineSMOTE",
            "Data-level balancing of the minority class",
            "Does not change the objective"],
           ["Devlin (2019);\nSanh (2019)", "BERT / DistilBERT",
            "Contextual embeddings for bug report text",
            "Needs task-specific fine-tuning"]],
          col_w=[2.35, 2.35, 4.05, 3.40], head_size=13.5, body_size=12.5,
          row_h=0.60)

    # ------------------------------------------------------------------
    # 6. RESEARCH GAPS
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Research Gaps",
                           "Three gaps identified; this study addresses all three")
    gaps = [
        ("GAP 1", "No unified comparison",
         "Traditional ML, hybrid deep models and transformers are rarely "
         "compared on the same data under one preprocessing pipeline and "
         "one evaluation protocol.", False),
        ("GAP 2", "Focal Loss unvalidated for defect prediction",
         "Algorithm-level correction is established in vision and NLP, but "
         "has not been validated for software defect prediction, least of "
         "all a class-weighted variant.", True),
        ("GAP 3", "Modalities studied separately",
         "Code-metric defect prediction and bug report classification are "
         "treated in disjoint literatures despite serving consecutive "
         "stages of the same process.", False),
    ]
    x = 0.60
    CARD_H = 4.05
    for tag, title, body, selected in gaps:
        w = 3.95
        fill = BRAND if selected else BAND
        rect(s, x, top + 0.10, w, CARD_H, fill=fill)
        tf = textbox(s, x + 0.30, top + 0.42, w - 0.60, CARD_H - 1.05)
        write(tf, tag, size=12, font=F_SUB,
              color=RGBColor(0xFF, 0xFF, 0xFF) if selected else MUTED,
              first=True, space_after=8)
        write(tf, title, size=19, font=F_SUB,
              color=RGBColor(0xFF, 0xFF, 0xFF) if selected else BRAND,
              space_after=12, line=1.12)
        write(tf, body, size=14, font=F_BODY,
              color=RGBColor(0xE6, 0xEC, 0xF6) if selected else BODY_CLR,
              space_after=0, line=1.25)
        if selected:
            tf = textbox(s, x + 0.30, top + CARD_H - 0.52, w - 0.60, 0.34)
            write(tf, "▶  SELECTED FOR THIS STUDY", size=12, font=F_SUB,
                  color=RGBColor(0xFF, 0xFF, 0xFF), first=True, space_after=0)
        x += w + 0.28

    # ------------------------------------------------------------------
    # 7. RESEARCH OBJECTIVES
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Research Objectives",
                           "Overall aim, specific objectives and research questions")

    rect(s, 0.60, top, 12.15, 0.82, fill=BAND)
    rect(s, 0.60, top, 0.075, 0.82, fill=BRAND)
    tf = textbox(s, 0.92, top + 0.14, 11.6, 0.58)
    rich(tf, [("AIM   ", F_SUB, MUTED, True),
              ("To empirically evaluate traditional ML and deep learning for "
               "bug detection, and determine how far imbalance-aware training "
               "improves detection on severely imbalanced defect data.",
               F_BODY, BODY_CLR, False)],
         size=15, first=True, space_after=0, line=1.2)

    tf = textbox(s, 0.60, top + 1.06, 6.05, 3.4)
    write(tf, "SPECIFIC OBJECTIVES", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=11)
    for n, t in [
        ("O1", "Implement Random Forest and Naive Bayes baselines"),
        ("O2", "Design and evaluate a hybrid CNN-LSTM architecture"),
        ("O3", "Fine-tune DistilBERT for bug report severity"),
        ("O4", "Implement and validate imbalance-aware training "
               "(Focal Loss, Class-Weighted Focal Loss)"),
        ("O5", "Compare all models on one held-out test set and by "
               "10-fold cross-validation"),
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(9)
        p.line_spacing = 1.15
        hanging(p, 0.45)
        r = p.add_run()
        r.text = n + "   "
        r.font.size = Pt(15)
        r.font.name = F_SUB
        r.font.bold = True
        r.font.color.rgb = BRAND
        r = p.add_run()
        r.text = t
        r.font.size = Pt(15)
        r.font.name = F_BODY
        r.font.color.rgb = BODY_CLR

    rect(s, 6.90, top + 1.06, 0.02, 3.3, fill=RULE)

    tf = textbox(s, 7.25, top + 1.06, 5.5, 3.4)
    write(tf, "RESEARCH QUESTIONS", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=11)
    for n, t in [
        ("RQ1", "How effectively do traditional ML models predict defects "
                "compared with deep learning?"),
        ("RQ2", "How do CNN-LSTM and BERT improve detection on large, "
                "imbalanced datasets?"),
        ("RQ3", "What is the impact of class imbalance on performance and "
                "generalisation?"),
        ("RQ4", "Which model gives the best balance of accuracy, precision, "
                "recall and F1?"),
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(9)
        p.line_spacing = 1.15
        hanging(p, 0.45)
        r = p.add_run()
        r.text = n + "   "
        r.font.size = Pt(15)
        r.font.name = F_SUB
        r.font.bold = True
        r.font.color.rgb = BRAND
        r = p.add_run()
        r.text = t
        r.font.size = Pt(15)
        r.font.name = F_BODY
        r.font.color.rgb = BODY_CLR

    # ------------------------------------------------------------------
    # 8. METHODOLOGY — design, data, preprocessing
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Methodology",
                           "Research design, data collection and preprocessing")

    tf = textbox(s, 0.60, top, 5.85, 1.05)
    write(tf, "DESIGN", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=7)
    write(tf, "Quantitative, experimental and comparative. Seven "
              "configurations vary model capacity and imbalance correction; "
              "everything else is held constant. All seeds fixed at 42.",
          size=15, font=F_BODY, color=BODY_CLR, space_after=0, line=1.2)

    tf = textbox(s, 0.60, top + 1.34, 5.85, 0.3)
    write(tf, "DATA", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=0)
    table(s, 0.60, top + 1.67, 5.85, 1.9,
          [["Property", "JM1 dataset"],
           ["Source", "PROMISE (NASA), via OpenML"],
           ["Modules", "10,885 C modules"],
           ["Features", "21 Halstead & McCabe metrics"],
           ["Defective", "2,106  (19.3%)"],
           ["*Imbalance", "*4.2 : 1"]],
          col_w=[2.15, 3.70], head_size=12, body_size=12)

    tf = textbox(s, 0.60, top + 3.76, 5.85, 0.5)
    write(tf, "Plus 1,200 bug reports (600 high / 600 low priority) for the "
              "classification task.", size=13, font=F_BODY, color=MUTED,
          first=True, space_after=0, line=1.2)

    tf = textbox(s, 6.85, top, 5.9, 0.3)
    write(tf, "PREPROCESSING PIPELINE", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=0)
    stages = [
        ("1", "Median imputation", "25 missing values, 5 columns"),
        ("2", "Winsorisation", "capped at the 99th percentile"),
        ("3", "Stratified split", "8,708 train / 2,177 test"),
        ("4", "Min-max normalisation", "fitted on training data only"),
        ("5", "BorderlineSMOTE", "train only → 14,046 balanced"),
    ]
    yy = top + 0.38
    for num, name, detail in stages:
        rect(s, 6.85, yy, 0.42, 0.42, fill=BRAND)
        tf = textbox(s, 6.85, yy + 0.09, 0.42, 0.26)
        write(tf, num, size=13, font=F_SUB, color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER, first=True, space_after=0)
        tf = textbox(s, 7.45, yy + 0.01, 5.3, 0.42)
        write(tf, name, size=15, font=F_SUB, color=BRAND, first=True,
              space_after=1)
        write(tf, detail, size=12.5, font=F_BODY, color=MUTED, space_after=0)
        yy += 0.63

    tf = textbox(s, 6.85, top + 3.62, 5.9, 0.5)
    write(tf, "The test set keeps its natural 4.2:1 imbalance, so results "
              "reflect deployment conditions.", size=13, font=F_BODY,
          color=MUTED, italic=True, first=True, space_after=0, line=1.2)

    # ------------------------------------------------------------------
    # 9. METHODOLOGY — models
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Methodology",
                           "Model families and the hybrid CNN-LSTM architecture")

    tf = textbox(s, 0.60, top, 7.0, 0.3)
    write(tf, "MODELS EVALUATED", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=0)
    table(s, 0.60, top + 0.32, 7.0, 2.55,
          [["Model", "Type", "Imbalance strategy"],
           ["Random Forest", "Traditional ML", "None / BorderlineSMOTE"],
           ["Gaussian Naive Bayes", "Traditional ML", "None / BorderlineSMOTE"],
           ["CNN-LSTM (BCE)", "Deep learning", "None"],
           ["CNN-LSTM (FL)", "Deep learning", "Focal Loss"],
           ["*CNN-LSTM (CW-FL)", "*Deep learning", "*SMOTE + Class-Weighted FL"],
           ["DistilBERT", "Transformer", "Balanced dataset"]],
          col_w=[2.65, 1.90, 2.45], head_size=12, body_size=12)

    tf = textbox(s, 0.60, top + 3.10, 7.0, 0.9)
    write(tf, "Each metric vector is reshaped to a 21-step sequence: "
              "convolutions detect local patterns among related metrics, "
              "the LSTM aggregates them across the vector.",
          size=13.5, font=F_BODY, color=MUTED, first=True, space_after=0,
          line=1.22)

    tf = textbox(s, 8.05, top, 4.7, 0.3)
    write(tf, "CNN-LSTM ARCHITECTURE", size=12, font=F_SUB, color=MUTED,
          first=True, space_after=0)
    layers = [
        ("Input", "(batch, 21, 1)"),
        ("Conv1D + BatchNorm", "64 filters, kernel 3"),
        ("Conv1D + BatchNorm", "128 filters, kernel 3"),
        ("LSTM", "64 hidden units"),
        ("Dropout", "rate 0.4"),
        ("Dense", "32 units, ReLU"),
        ("Output", "1 unit, sigmoid"),
    ]
    yy = top + 0.36
    for i, (name, cfg) in enumerate(layers):
        shade = BRAND if i in (1, 2, 3) else BAND
        rect(s, 8.05, yy, 4.7, 0.42, fill=shade)
        tf = textbox(s, 8.25, yy + 0.09, 2.6, 0.26)
        write(tf, name, size=13, font=F_SUB,
              color=RGBColor(0xFF, 0xFF, 0xFF) if shade == BRAND else BRAND,
              first=True, space_after=0)
        tf = textbox(s, 10.85, yy + 0.10, 1.75, 0.26)
        write(tf, cfg, size=11.5, font=F_BODY,
              color=RGBColor(0xD8, 0xE0, 0xF0) if shade == BRAND else MUTED,
              align=PP_ALIGN.RIGHT, first=True, space_after=0)
        yy += 0.50

    tf = textbox(s, 8.05, yy + 0.04, 4.7, 0.3)
    write(tf, "77,121 trainable parameters", size=13, font=F_SUB,
          color=BRAND, align=PP_ALIGN.RIGHT, first=True, space_after=0)

    # ------------------------------------------------------------------
    # 10. PROPOSED METHOD
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Proposed Method",
                           "Class-Weighted Focal Loss — the contribution "
                           "of this study")

    y = top + 0.05
    for label, formula, params, note, hi in [
        ("Binary Cross-Entropy",
         "BCE  =  −[ y·log(p) + (1−y)·log(1−p) ]",
         None,
         "Every example counts equally, so the majority class dominates "
         "the gradient.", False),
        ("Focal Loss  (Lin et al., 2017)",
         "FL  =  −α · (1 − pₜ)^γ · log(pₜ)",
         "α = 0.75      γ = 2.0",
         "Example-level: suppresses easy examples so training focuses on "
         "hard ones.", False),
        ("Class-Weighted Focal Loss  (proposed)",
         "CW-FL  =  w(y) · FL",
         "w = 4.2 if defective      w = 1.0 if clean",
         "Class-level: the observed imbalance ratio is encoded into the "
         "objective itself.", True),
    ]:
        h = 1.45
        white = RGBColor(0xFF, 0xFF, 0xFF)
        rect(s, 0.60, y, 7.55, h, fill=BRAND if hi else BAND)
        if not hi:
            rect(s, 0.60, y, 0.06, h, fill=BRAND)
        tf = textbox(s, 0.95, y + 0.14, 6.95, 1.20)
        write(tf, label, size=14, font=F_SUB,
              color=white if hi else BRAND, first=True, space_after=5)
        write(tf, formula, size=15, font="Consolas",
              color=white if hi else BODY_CLR, space_after=3)
        if params:
            write(tf, params, size=12.5, font="Consolas",
                  color=RGBColor(0xC8, 0xD4, 0xEA) if hi else MUTED,
                  space_after=5)
        write(tf, note, size=12.5, font=F_BODY,
              color=RGBColor(0xD8, 0xE0, 0xF0) if hi else MUTED,
              space_after=0, line=1.16)
        y += h + 0.18

    rect(s, 8.50, top + 0.05, 4.28, 4.71, fill=BAND)
    tf = textbox(s, 8.82, top + 0.34, 3.65, 4.1)
    write(tf, "WHY IT WORKS", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=12)
    for t in [
        "The focusing term is adaptive: it tracks how well the model "
        "currently classifies each example.",
        "The class weight is static: it encodes a fixed property of the "
        "dataset.",
        "They correct different aspects of the same problem, so their "
        "effects add.",
    ]:
        bullet(tf, t, size=13.5, space_after=11)
    write(tf, "SIGNIFICANCE", size=12, font=F_SUB, color=MUTED,
          space_before=8, space_after=8)
    write(tf, "First systematic validation of class-weighted Focal Loss on "
              "a standard software defect prediction benchmark.",
          size=13.5, font=F_BODY, color=BRAND, space_after=0, line=1.2)

    # ------------------------------------------------------------------
    # 11. RESULTS — baselines
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Results",
                           "Baseline models on the held-out test set "
                           "(2,177 modules, 421 defective)")
    table(s, 0.60, top + 0.05, 7.4, 2.2,
          [["Model", "Acc.", "Prec.", "Recall", "F1", "Detected"],
           ["RF, no SMOTE", "0.811", "0.528", "0.202", "0.292", "85 / 421"],
           ["NB, no SMOTE", "0.790", "0.441", "0.321", "0.371", "135 / 421"],
           ["NB + SMOTE", "0.774", "0.404", "0.359", "0.380", "151 / 421"],
           ["*RF + SMOTE", "*0.765", "*0.404", "*0.451", "*0.427",
            "*190 / 421"]],
          col_w=[2.25, 0.95, 0.95, 1.00, 0.95, 1.30], head_size=12,
          body_size=12.5)

    tf = textbox(s, 0.60, top + 2.35, 7.4, 1.9)
    for t in [
        "Highest accuracy (RF, no SMOTE) finds the fewest defects — "
        "the accuracy paradox in one row.",
        "BorderlineSMOTE raises RF recall by 123%: 85 → 190 defects "
        "detected, at 4.5 points of accuracy.",
        "SMOTE barely moves Naive Bayes, which fits class-conditional "
        "Gaussians rather than a boundary.",
    ]:
        bullet(tf, t, size=15, space_after=11)

    rect(s, 8.40, top + 0.05, 4.38, 1.72, fill=BAND)
    rect(s, 8.40, top + 0.05, 0.075, 1.72, fill=BRAND)
    tf = textbox(s, 8.72, top + 0.28, 3.85, 1.35)
    write(tf, "BEST TRADITIONAL BASELINE", size=11.5, font=F_SUB,
          color=MUTED, first=True, space_after=6)
    write(tf, "Random Forest + SMOTE", size=17, font=F_SUB, color=BRAND,
          space_after=5)
    write(tf, "Recall 0.451  ·  190 of 421 found  ·  231 missed",
          size=13, font=F_BODY, color=BODY_CLR, space_after=0, line=1.2)

    picture(s, "11_baseline_comparison.png", 8.40, top + 1.95, w=4.38)

    # ------------------------------------------------------------------
    # 12. RESULTS — effect of the proposed loss
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Results",
                           "Effect of the loss function — architecture, "
                           "data and protocol held constant")
    table(s, 0.60, top + 0.05, 7.4, 1.85,
          [["CNN-LSTM configuration", "Acc.", "Prec.", "Recall", "F1",
            "Detected"],
           ["Binary Cross-Entropy", "0.665", "0.306", "0.577", "0.400",
            "243 / 421"],
           ["Focal Loss", "0.558", "0.270", "0.755", "0.398", "318 / 421"],
           ["*Class-Weighted Focal Loss", "*0.461", "*0.241", "*0.834",
            "*0.374", "*351 / 421"]],
          col_w=[2.75, 0.85, 0.85, 0.95, 0.85, 1.15], head_size=12,
          body_size=12.5)

    tf = textbox(s, 0.60, top + 2.12, 7.4, 2.3)
    for t, c in [
        ("Even with plain BCE the CNN-LSTM beats every traditional model "
         "on recall — the architecture contributes independently of "
         "the loss.", BODY_CLR),
        ("Focal Loss adds 17.8 points of recall and lifts ROC-AUC from "
         "0.681 to 0.700, so it improved the learned ranking, not just "
         "the threshold.", BODY_CLR),
        ("The class weight adds a further 7.9 points, reaching "
         "recall 0.834.", BRAND),
        ("Every gain in recall costs precision. This is a trade-off, "
         "not a free lunch.", MUTED),
    ]:
        bullet(tf, t, size=15, color=c, space_after=11)

    picture(s, "21_bugs_caught.png", 8.30, top + 0.30, w=4.5)

    # ------------------------------------------------------------------
    # 13. RESULTS — overall comparison + DistilBERT
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Results",
                           "Overall comparison and bug report classification")

    picture(s, "19_final_comparison_bar.png", 0.60, top + 0.05, w=7.5)

    rect(s, 8.40, top + 0.05, 4.38, 1.80, fill=BRAND)
    tf = textbox(s, 8.72, top + 0.24, 3.80, 1.45)
    write(tf, "HEADLINE RESULT", size=11.5, font=F_SUB,
          color=RGBColor(0xC8, 0xD4, 0xEA), first=True, space_after=4)
    write(tf, "+84.7%", size=32, font=F_SUB, color=RGBColor(0xFF, 0xFF, 0xFF),
          space_after=4)
    write(tf, "detection rate over the best traditional baseline "
              "(190 → 351 of 421 defects)", size=12.5, font=F_BODY,
          color=RGBColor(0xD8, 0xE0, 0xF0), space_after=0, line=1.18)

    tf = textbox(s, 8.40, top + 2.00, 4.38, 0.3)
    write(tf, "RECALL, CODE-METRIC MODELS", size=11.5, font=F_SUB,
          color=MUTED, first=True, space_after=0)

    yy = top + 2.36
    for label, val, frac, c in [("RF + SMOTE", "0.451", 0.451, RULE),
                                ("CNN-LSTM CW-FL", "0.834", 0.834, BRAND)]:
        tf = textbox(s, 8.40, yy, 4.38, 0.24)
        write(tf, label, size=12.5, font=F_BODY, color=BODY_CLR, first=True,
              space_after=0)
        tf = textbox(s, 8.40, yy, 4.38, 0.24)
        write(tf, val, size=12.5, font=F_SUB, color=BRAND,
              align=PP_ALIGN.RIGHT, first=True, space_after=0)
        rect(s, 8.40, yy + 0.26, 4.38, 0.16, fill=RGBColor(0xE8, 0xEC, 0xF2))
        rect(s, 8.40, yy + 0.26, 4.38 * frac, 0.16, fill=c)
        yy += 0.62

    rect(s, 8.40, top + 3.72, 4.38, 1.10, fill=BAND)
    tf = textbox(s, 8.72, top + 3.92, 3.80, 0.85)
    write(tf, "DISTILBERT — BUG REPORTS", size=11.5, font=F_SUB,
          color=MUTED, first=True, space_after=5)
    write(tf, "F1 = 1.000 on 240 held-out reports. The classes are "
              "lexically separable; expect 0.85–0.95 on live tracker "
              "data.", size=12.5, font=F_BODY, color=BODY_CLR, space_after=0,
          line=1.18)

    # ------------------------------------------------------------------
    # 14. DISCUSSION
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Discussion",
                           "What the results mean, and what they cost")

    cards = [
        ("The trade-off is the finding",
         "CW-FL flags 1,455 of 2,177 modules and finds 351 defects; "
         "RF+SMOTE flags 470 and finds 190. Yield per review falls from "
         "0.40 to 0.24, but coverage reaches a level RF cannot attain at "
         "any threshold."),
        ("F1 misranks these models",
         "F1 peaks at RF+SMOTE and would reject the model that finds 161 "
         "more defects. F1 weights both errors equally; a defect reaching "
         "production costs ~10x a wasted review."),
        ("Architecture and objective both matter",
         "BCE-trained CNN-LSTM already beats RF+SMOTE on recall, so the "
         "gain is not from the loss alone. The factorial design separates "
         "the two contributions."),
    ]
    x = 0.60
    for title, body in cards:
        rect(s, x, top + 0.10, 3.95, 2.35, fill=BAND)
        rect(s, x, top + 0.10, 3.95, 0.06, fill=BRAND)
        tf = textbox(s, x + 0.28, top + 0.42, 3.40, 1.9)
        write(tf, title, size=16.5, font=F_SUB, color=BRAND, first=True,
              space_after=9, line=1.12)
        write(tf, body, size=13, font=F_BODY, color=BODY_CLR, space_after=0,
              line=1.22)
        x += 4.23

    rect(s, 0.60, top + 2.68, 12.15, 1.62, fill=BAND)
    rect(s, 0.60, top + 2.68, 0.075, 1.62, fill=ACCENT)
    tf = textbox(s, 0.92, top + 2.86, 11.6, 1.3)
    write(tf, "LIMITATIONS", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=8)
    write(tf, "Single dataset (JM1), no cross-project validation  ·  "
              "constructed rather than live bug report data  ·  "
              "no cross-validation of the deep models  ·  "
              "CPU-only training limited hyperparameter search  ·  "
              "static features only, no process metrics  ·  "
              "no explainability analysis",
          size=14, font=F_BODY, color=BODY_CLR, space_after=0, line=1.3)

    # ------------------------------------------------------------------
    # 15. FUTURE WORKS
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Future Works",
                           "Ordered by proximity to the present study")
    items = [
        ("01", "Cross-project validation",
         "Extend to KC1 and PC1, then train on one project and test on "
         "another to test generalisation under distribution shift."),
        ("02", "Search the loss hyperparameters",
         "α, γ and ρ were set on principled grounds. A "
         "systematic search would test whether ρ should equal the "
         "imbalance ratio."),
        ("03", "Cost-sensitive evaluation",
         "Replace F1 with Fβ calibrated to a real cost ratio, or "
         "expected-cost curves, so selection stops relying on judgement."),
        ("04", "Live bug report data",
         "Reproduce the DistilBERT result on Eclipse or Mozilla Bugzilla "
         "exports with authentic label noise."),
        ("05", "Explainability",
         "Apply SHAP or LIME so a flagged module comes with a reason. "
         "A model that cannot explain itself will not be trusted."),
        ("06", "Richer features",
         "Add process metrics (churn, defect history) and apply graph "
         "neural networks to abstract syntax trees."),
    ]
    x, y = 0.60, top + 0.10
    for i, (num, title, body) in enumerate(items):
        col, row = i % 3, i // 3
        xx = x + col * 4.23
        yy = y + row * 2.12
        rect(s, xx, yy, 3.95, 1.92, fill=BAND)
        tf = textbox(s, xx + 0.28, yy + 0.22, 3.4, 0.3)
        write(tf, num, size=13, font=F_SUB, color=ACCENT, first=True,
              space_after=0)
        tf = textbox(s, xx + 0.28, yy + 0.56, 3.4, 1.2)
        write(tf, title, size=15.5, font=F_SUB, color=BRAND, first=True,
              space_after=7, line=1.1)
        write(tf, body, size=12.5, font=F_BODY, color=BODY_CLR, space_after=0,
              line=1.2)

    # ------------------------------------------------------------------
    # 16. PUBLICATIONS & REFERENCES
    # ------------------------------------------------------------------
    s, top = content_slide(prs, "Publications & References", None)

    rect(s, 0.60, top, 12.15, 0.78, fill=BAND)
    rect(s, 0.60, top, 0.075, 0.78, fill=BRAND)
    tf = textbox(s, 0.92, top + 0.16, 11.6, 0.55)
    rich(tf, [("PUBLICATIONS   ", F_SUB, MUTED, True),
              ("None at the time of submission. A manuscript reporting the "
               "Class-Weighted Focal Loss results is in preparation for a "
               "peer-reviewed software engineering conference.",
               F_BODY, BODY_CLR, False)],
         size=14, first=True, space_after=0, line=1.2)

    tf = textbox(s, 0.60, top + 1.02, 6.05, 3.3)
    write(tf, "KEY REFERENCES", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=10)
    left = [
        "[1] B. W. Boehm, Software Engineering Economics. Prentice-Hall, "
        "1981.",
        "[2] L. Breiman, “Random forests,” Machine Learning, "
        "vol. 45, no. 1, pp. 5–32, 2001.",
        "[3] N. V. Chawla et al., “SMOTE,” JAIR, vol. 16, "
        "pp. 321–357, 2002.",
        "[4] J. Devlin et al., “BERT,” in Proc. NAACL-HLT, 2019, "
        "pp. 4171–4186.",
        "[5] M. H. Halstead, Elements of Software Science. Elsevier, 1977.",
        "[6] H. Han et al., “Borderline-SMOTE,” in Proc. ICIC, "
        "2005, pp. 878–887.",
    ]
    for t in left:
        write(tf, t, size=12.5, font=F_BODY, color=BODY_CLR, space_after=9,
              line=1.18)

    rect(s, 6.90, top + 1.02, 0.02, 3.2, fill=RULE)

    tf = textbox(s, 7.25, top + 1.02, 5.5, 3.3)
    write(tf, " ", size=12, font=F_SUB, color=MUTED, first=True,
          space_after=10)
    right = [
        "[7] S. Hochreiter and J. Schmidhuber, “Long short-term "
        "memory,” Neural Computation, vol. 9, no. 8, 1997.",
        "[8] T.-Y. Lin et al., “Focal loss for dense object "
        "detection,” in Proc. ICCV, 2017, pp. 2999–3007.",
        "[9] T. J. McCabe, “A complexity measure,” IEEE TSE, "
        "vol. SE-2, no. 4, pp. 308–320, 1976.",
        "[10] T. Menzies et al., “Defect prediction from static code "
        "features,” ASE, vol. 17, no. 4, 2010.",
        "[11] V. Sanh et al., “DistilBERT,” arXiv:1910.01108, "
        "2019.",
        "[12] J. S. Shirabad and T. J. Menzies, The PROMISE Repository, "
        "Univ. of Ottawa, 2005.",
    ]
    for t in right:
        write(tf, t, size=12.5, font=F_BODY, color=BODY_CLR, space_after=9,
              line=1.18)

    # ------------------------------------------------------------------
    # 17. THANK YOU
    # ------------------------------------------------------------------
    s = new_slide(prs)
    chrome(s)
    tf = textbox(s, 1.60, 2.75, 10.1, 1.15, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, "Thank You", size=54, font=F_TITLE, color=BRAND,
          align=PP_ALIGN.CENTER, first=True, space_after=0)
    rect(s, 6.17, 4.06, 1.0, 0.04, fill=BRAND)
    tf = textbox(s, 1.60, 4.30, 10.1, 0.9)
    write(tf, "Questions & Discussion", size=21, font=F_SUB, color=MUTED,
          align=PP_ALIGN.CENTER, first=True, space_after=8)
    write(tf, "S. Shanthosh   ·   20APSE4882   ·   "
              "researchp180@gmail.com", size=14, font=F_BODY, color=BODY_CLR,
          align=PP_ALIGN.CENTER, space_after=0)

    prs.save(OUTPUT)
    print("Saved:", OUTPUT)
    print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    build()
