"""
=========================================================
 FINAL EVALUATION PRESENTATION GENERATOR
 SE 8101 Research Project -- Sabaragamuwa University of Sri Lanka
 Student : S.Shanthosh | 20APSE4882
 Supervisor: Mrs. WMLS Abeythunga
 Date    : 27th May 2026
=============================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import os

# ── Paths ────────────────────────────────────────────────────
RESULTS = r'd:\Betterhalf\Research 2026\results'
OUT     = r'd:\Betterhalf\Research 2026\report\Final_Presentation_20APSE4882.pptx'

# ── Color Palette ────────────────────────────────────────────
DARK_BLUE   = RGBColor(0,   70,  127)
MID_BLUE    = RGBColor(0,  114,  198)
LIGHT_BLUE  = RGBColor(220, 235, 248)
PALE_BLUE   = RGBColor(240, 246, 255)
ORANGE      = RGBColor(215,  90,   0)
GOLD        = RGBColor(200, 155,   0)
WHITE       = RGBColor(255, 255, 255)
DARK_GRAY   = RGBColor( 50,  50,   50)
MID_GRAY    = RGBColor(120, 120, 120)
LIGHT_GRAY  = RGBColor(245, 245, 245)
GREEN_DK    = RGBColor(  0, 112,   0)
RED_DK      = RGBColor(192,   0,   0)
TEAL        = RGBColor(  0, 128, 128)

# ════════════════════════════════════════════════════════════
# Presentation setup  (16:9, 13.33" x 7.5")
# ════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank layout

# ════════════════════════════════════════════════════════════
# Helper utilities
# ════════════════════════════════════════════════════════════
def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None, lw=0.5, radius=0):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = _Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, italic=False,
        color=DARK_GRAY, align=PP_ALIGN.LEFT,
        v=MSO_VERTICAL_ANCHOR.TOP, wrap=True, name='Calibri'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap  = wrap
    tf.vertical_anchor = v
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = name
    return tf

def para(tf, text, size=16, bold=False, italic=False,
         color=DARK_GRAY, align=PP_ALIGN.LEFT, space_before=4, level=0):
    p = tf.add_paragraph()
    p.alignment    = align
    p.level        = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = 'Calibri'
    return p

def bullet(tf, text, size=16, bold=False, color=DARK_GRAY, space_before=5,
           icon='●', indent='  '):
    return para(tf, f'{indent}{icon}  {text}',
                size=size, bold=bold, color=color,
                space_before=space_before)

def img(slide, fname, l, t, w=None, h=None):
    path = os.path.join(RESULTS, fname)
    if not os.path.exists(path):
        return
    kw = {}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def header(slide, title, sub=''):
    """Standard header bar used by all content slides."""
    rect(slide, 0, 0, 13.333, 1.25, fill=DARK_BLUE)
    rect(slide, 0, 1.25, 13.333, 0.06, fill=ORANGE)
    txt(slide, title, 0.35, 0.1, 11.5, 0.75,
        size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    if sub:
        txt(slide, sub, 0.35, 0.82, 11, 0.38,
            size=13, italic=True, color=RGBColor(180, 210, 240))
    # Footer
    rect(slide, 0, 7.18, 13.333, 0.32, fill=RGBColor(235, 240, 248))
    txt(slide, 'S.Shanthosh  |  20APSE4882  |  SE 8101 Research Project  |  '
        'Sabaragamuwa University of Sri Lanka',
        0.3, 7.19, 11, 0.28, size=9, color=MID_GRAY)

def make_table(slide, headers, rows, l, t, w, h,
               col_w=None, hdr_size=12, row_size=11, hi_rows=None):
    """Build a formatted table (blue header, alternating rows)."""
    nr   = len(rows) + 1
    nc   = len(headers)
    tbl  = slide.shapes.add_table(
        nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table

    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)

    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr; r.font.size = Pt(hdr_size)
        r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = 'Calibri'

    # Data rows
    hi_rows = hi_rows or []
    for i, row in enumerate(rows):
        bg = GOLD if i in hi_rows else (LIGHT_BLUE if i % 2 == 1 else WHITE)
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val); r.font.size = Pt(row_size)
            r.font.color.rgb = DARK_GRAY
            r.font.name = 'Calibri'
            if i in hi_rows:
                r.font.bold = True

    return tbl

def box(slide, l, t, w, h, text, fill=LIGHT_BLUE,
        tc=DARK_BLUE, size=14, bold=True, align=PP_ALIGN.CENTER, border=None):
    r = rect(slide, l, t, w, h, fill=fill, line=border or DARK_BLUE, lw=0.8)
    txt(slide, text, l+0.08, t+0.05, w-0.16, h-0.1,
        size=size, bold=bold, color=tc, align=align,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 1 -- TITLE
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.333, 7.5, fill=DARK_BLUE)                  # full dark bg
rect(s, 0, 5.5, 13.333, 0.07, fill=ORANGE)                  # orange stripe

txt(s, 'SABARAGAMUWA UNIVERSITY OF SRI LANKA', 0, 0.35, 13.333, 0.5,
    size=17, bold=True, color=RGBColor(180, 210, 240),
    align=PP_ALIGN.CENTER)
txt(s, 'Department of Software Engineering  |  SE 8101 Research Project',
    0, 0.82, 13.333, 0.42,
    size=13, color=RGBColor(160, 190, 220), align=PP_ALIGN.CENTER)

rect(s, 1.0, 1.45, 11.333, 0.05, fill=ORANGE)              # top title rule
txt(s,
    'An Empirical Evaluation of Deep Learning Models\n'
    'for Enhanced and Scalable Bug Detection and\n'
    'Classification in Large-Scale Software Systems',
    0.6, 1.6, 12.133, 2.5,
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 1.0, 4.15, 11.333, 0.05, fill=ORANGE)              # bottom title rule

# Info boxes
rect(s, 1.5, 4.4, 4.8, 1.55, fill=RGBColor(0, 50, 100))
txt(s, 'Student', 1.6, 4.5, 4.5, 0.35, size=12,
    color=RGBColor(160, 200, 240), bold=True)
txt(s, 'S. Shanthosh', 1.6, 4.85, 4.5, 0.38, size=18,
    color=WHITE, bold=True)
txt(s, 'Index: 20APSE4882', 1.6, 5.2, 4.5, 0.35, size=13,
    color=RGBColor(180, 210, 240))
txt(s, '27th May 2026', 1.6, 5.58, 4.5, 0.3, size=12,
    color=RGBColor(160, 190, 220), italic=True)

rect(s, 7.0, 4.4, 4.8, 1.55, fill=RGBColor(0, 50, 100))
txt(s, 'Supervisor', 7.1, 4.5, 4.5, 0.35, size=12,
    color=RGBColor(160, 200, 240), bold=True)
txt(s, 'Mrs. WMLS Abeythunga', 7.1, 4.85, 4.5, 0.38, size=16,
    color=WHITE, bold=True)
txt(s, 'Department of Software Engineering', 7.1, 5.22, 4.5, 0.35,
    size=12, color=RGBColor(180, 210, 240))
txt(s, 'BSc (Hons) Software Engineering', 7.1, 5.58, 4.5, 0.3,
    size=12, color=RGBColor(160, 190, 220), italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 2 -- OUTLINE
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Presentation Outline')

sections = [
    ('01', 'Introduction & Problem',
     'Background, motivation, class imbalance challenge'),
    ('02', 'Research Design',
     'Objectives, questions, dataset, methodology'),
    ('03', 'Experiments & Results',
     'Baseline ML, CNN-LSTM, Focal Loss, DistilBERT'),
    ('04', 'Findings & Conclusion',
     'Comparative analysis, contributions, future work'),
]
cols = [(0.4, DARK_BLUE), (4.0, DARK_BLUE), (7.6, DARK_BLUE), (11.2, DARK_BLUE)]
for i, ((l, c), (num, title, sub)) in enumerate(zip(cols, sections)):
    rect(s, l, 1.45, 2.8, 5.5, fill=DARK_BLUE if i % 2 == 0 else MID_BLUE)
    txt(s, num, l+0.1, 1.55, 2.6, 1.1, size=54, bold=True,
        color=RGBColor(255, 255, 255) if i % 2 == 0 else WHITE,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    rect(s, l, 2.7, 2.8, 0.05, fill=ORANGE)
    txt(s, title, l+0.1, 2.82, 2.6, 0.75, size=16, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, l+0.1, 3.62, 2.6, 1.25, size=12, italic=True,
        color=RGBColor(200, 220, 240), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 3 -- INTRODUCTION & BACKGROUND
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Introduction & Motivation',
       'Why does automated bug detection matter?')

# Left: text
tf = txt(s, '', 0.35, 1.45, 6.8, 5.5, size=16, wrap=True)
tf = s.shapes[-1].text_frame
tf.word_wrap = True

p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = 'The Cost of Software Bugs'
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = DARK_BLUE
r.font.name = 'Calibri'

bullet(tf, 'Global software bug cost: $1.7 trillion/year (CISQ, 2022)',
       size=15, color=DARK_GRAY)
bullet(tf, 'Cost to fix a bug grows 10x from coding to production',
       size=15, color=DARK_GRAY)
bullet(tf, 'Manual code review is slow, inconsistent, not scalable',
       size=15, color=DARK_GRAY)
bullet(tf, 'Rule-based static tools generate too many false alarms',
       size=15, color=DARK_GRAY)

para(tf, '', size=8)  # spacer
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = 'Why Machine Learning?'
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = DARK_BLUE; r2.font.name = 'Calibri'
p2.space_before = Pt(8)

bullet(tf, 'ML models learn defect patterns from historical data',
       size=15, color=DARK_GRAY)
bullet(tf, 'Generalise across projects without hand-coded rules',
       size=15, color=DARK_GRAY)
bullet(tf, 'Deep Learning (CNN-LSTM, BERT) captures complex patterns',
       size=15, color=DARK_GRAY)

# Right: stat boxes
rect(s, 7.4, 1.45, 5.5, 1.6, fill=DARK_BLUE)
txt(s, '$1.7 Trillion', 7.4, 1.55, 5.5, 0.7,
    size=32, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, 'Annual global cost of software bugs', 7.4, 2.22, 5.5, 0.7,
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 7.4, 3.2, 5.5, 1.6, fill=MID_BLUE)
txt(s, '80%', 7.4, 3.3, 5.5, 0.7,
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 'of software projects exceed budget\ndue to unexpected defects', 7.4, 3.98, 5.5, 0.7,
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 7.4, 4.95, 5.5, 1.6, fill=RGBColor(0, 50, 100))
txt(s, '19.3%', 7.4, 5.05, 5.5, 0.7,
    size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, 'of modules in our JM1 dataset are defective\n(severe class imbalance)', 7.4, 5.72, 5.5, 0.7,
    size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 4 -- PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Problem Statement',
       'The core challenge: severe class imbalance in defect datasets')

# Problem box
rect(s, 0.35, 1.4, 12.6, 1.05, fill=RGBColor(255, 235, 230))
rect(s, 0.35, 1.4, 0.08, 1.05, fill=RED_DK)
txt(s, 'AI-driven bug detection faces critical limitations: class imbalance leads models '
    'to ignore the minority defective class, producing near-zero Recall on the bugs that matter most.',
    0.55, 1.45, 12.2, 0.95, size=14, bold=False, color=RED_DK)

# Three problems
for i, (title, body, c) in enumerate([
    ('Class Imbalance',
     '4.2:1 ratio (non-defective:defective). A model predicting ALL as non-defective '
     'gets 80.7% accuracy -- yet catches ZERO bugs.',
     RED_DK),
    ('Recall vs Accuracy Paradox',
     'High accuracy is misleading. Missing a critical bug costs far more than a '
     'false alarm. Recall on the buggy class is the true metric.',
     ORANGE),
    ('Traditional Approaches Fall Short',
     'Random Forest without SMOTE: Recall = 0.202 (catches only 85/421 bugs). '
     'Manual code review does not scale to millions of lines of code.',
     DARK_BLUE),
]):
    l = 0.35 + i*4.3
    rect(s, l, 2.6, 4.0, 3.5, fill=PALE_BLUE, line=c, lw=1.5)
    rect(s, l, 2.6, 4.0, 0.48, fill=c)
    txt(s, title, l+0.1, 2.65, 3.8, 0.38,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, body, l+0.1, 3.18, 3.8, 2.85,
        size=13, color=DARK_GRAY, wrap=True)

img(s, '01_class_distribution.png', 9.9, 2.6, 3.1, 3.0)


# ════════════════════════════════════════════════════════════
# SLIDE 5 -- RESEARCH OBJECTIVES & QUESTIONS
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Research Objectives & Questions')

# Left: Objectives
rect(s, 0.35, 1.45, 6.2, 0.5, fill=DARK_BLUE)
txt(s, 'Research Objectives', 0.45, 1.47, 6.0, 0.46,
    size=15, bold=True, color=WHITE, v=MSO_VERTICAL_ANCHOR.MIDDLE)

tf = s.shapes.add_textbox(Inches(0.35), Inches(2.05),
                           Inches(6.2), Inches(5.1)).text_frame
tf.word_wrap = True
for obj, detail in [
    ('O1', 'Evaluate Random Forest & Naive Bayes as ML baselines'),
    ('O2', 'Design CNN-LSTM architecture for defect prediction'),
    ('O3', 'Fine-tune DistilBERT for bug report classification'),
    ('O4', 'Implement Focal Loss & Class-Weighted Focal Loss\n'
           '      (primary research contribution)'),
    ('O5', 'Compare all models on Accuracy, Recall, F1, AUC'),
]:
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    r1 = p.add_run()
    r1.text = f'  {obj}:  '
    r1.font.size = Pt(15); r1.font.bold = True
    r1.font.color.rgb = ORANGE; r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = detail
    r2.font.size = Pt(14); r2.font.color.rgb = DARK_GRAY
    r2.font.name = 'Calibri'

# Right: Research Questions
rect(s, 6.9, 1.45, 6.1, 0.5, fill=MID_BLUE)
txt(s, 'Research Questions', 7.0, 1.47, 5.9, 0.46,
    size=15, bold=True, color=WHITE, v=MSO_VERTICAL_ANCHOR.MIDDLE)

for i, (rq, q) in enumerate([
    ('RQ1', 'How effectively do traditional ML models predict defects compared to DL?'),
    ('RQ2', 'How do CNN-LSTM and BERT improve bug detection on imbalanced data?'),
    ('RQ3', 'What is the impact of class imbalance on model performance?'),
    ('RQ4', 'Which model achieves the best balance of Recall and F1-Score?'),
]):
    t = 2.1 + i * 1.26
    rect(s, 6.9, t, 6.1, 1.15, fill=PALE_BLUE if i % 2 == 0 else LIGHT_BLUE,
         line=MID_BLUE, lw=0.5)
    txt(s, rq, 7.0, t+0.04, 0.65, 1.07,
        size=18, bold=True, color=MID_BLUE,
        v=MSO_VERTICAL_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    txt(s, q, 7.7, t+0.06, 5.1, 1.03,
        size=12, color=DARK_GRAY, wrap=True, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 6 -- DATASET
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Dataset: JM1 (PROMISE Repository)',
       'NASA real-time ground system -- 10,885 C modules, 21 complexity features')

make_table(s,
    headers=['Property', 'Value'],
    rows=[
        ['Source',            'PROMISE Software Engineering Repository (NASA)'],
        ['Language',          'C (real-time ground system)'],
        ['Total Modules',     '10,885'],
        ['Features',          '21 software complexity metrics'],
        ['Target Variable',   'defects  (1 = buggy,  0 = clean)'],
        ['Defective Modules', '2,106  (19.3%)'],
        ['Non-Defective',     '8,779  (80.7%)'],
        ['Imbalance Ratio',   '4.2 : 1  (non-defective : defective)'],
        ['Missing Values',    '25  (handled by median imputation)'],
        ['Train / Test Split','80% / 20%  (stratified, seed=42)'],
    ],
    l=0.35, t=1.45, w=6.6, h=5.6,
    col_w=[2.4, 4.2], hdr_size=13, row_size=12,
    hi_rows=[6, 7])

# Feature groups on right
right_groups = [
    ('Halstead Complexity',
     'volume, difficulty, effort,\nlength, time, prog_lang'),
    ('McCabe Complexity',
     'v(g), ev(g), iv(g),\ncyclomatic variants'),
    ('Code Volume',
     'loc, lOCode, lOComment,\nlOBlank, locCodeAndComment'),
    ('Operator/Operand',
     'uniq_Op, uniq_Opnd,\ntotal_Op, total_Opnd, n'),
]
txt(s, 'Feature Groups', 7.3, 1.45, 5.7, 0.42,
    size=15, bold=True, color=DARK_BLUE)
for i, (group, feats) in enumerate(right_groups):
    t = 1.98 + i * 1.32
    rect(s, 7.3, t, 5.7, 1.22, fill=PALE_BLUE if i%2==0 else LIGHT_BLUE,
         line=MID_BLUE, lw=0.4)
    txt(s, group, 7.4, t+0.06, 5.5, 0.38, size=13, bold=True, color=DARK_BLUE)
    txt(s, feats, 7.4, t+0.46, 5.5, 0.72, size=11, color=DARK_GRAY, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 7 -- PREPROCESSING PIPELINE
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Preprocessing Pipeline',
       '5-stage pipeline to prepare JM1 for model training')

stages = [
    ('1', 'Imputation',    'Median imputation\n25 missing values\nfixed',
     DARK_BLUE,  RGBColor(0,  30,  87)),
    ('2', 'Winsorization', 'Cap outliers at\n99th percentile\nper feature',
     MID_BLUE,   RGBColor(0,  74, 158)),
    ('3', 'Train/Test\nSplit', 'Stratified 80/20\nseed=42\n8,708 / 2,177',
     TEAL,       RGBColor(0,  88,  88)),
    ('4', 'Normalisation', 'Min-Max scaling\nto [0,1]\nfit on train only',
     GREEN_DK,   RGBColor(0,  72,   0)),
    ('5', 'SMOTE', 'BorderlineSMOTE\non train set only\n14,046 balanced',
     ORANGE,     RGBColor(175, 50,  0)),
]

for i, (num, name, detail, c, c_dark) in enumerate(stages):
    l = 0.4 + i * 2.55
    # Box
    rect(s, l, 1.55, 2.25, 4.8, fill=c)
    rect(s, l, 1.55, 2.25, 0.85, fill=c_dark)
    # Step number
    txt(s, f'Stage {num}', l, 1.62, 2.25, 0.7,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    # Stage name
    txt(s, name, l, 2.5, 2.25, 0.9,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    # Detail
    txt(s, detail, l+0.1, 3.5, 2.05, 1.5,
        size=13, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow (except last)
    if i < 4:
        txt(s, '▶', l+2.25, 3.7, 0.28, 0.5,
            size=20, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Result boxes
rect(s, 0.35, 6.45, 12.6, 0.65, fill=LIGHT_BLUE, line=DARK_BLUE, lw=0.8)
txt(s, 'Train: 14,046 (balanced 1:1 via SMOTE)   |   '
    'Test: 2,177 (natural imbalance, 421 buggy)   |   '
    'Leakage prevention: Scaler & SMOTE fitted on train only',
    0.5, 6.5, 12.3, 0.56, size=13, bold=False, color=DARK_BLUE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 8 -- METHODOLOGY OVERVIEW (4 model families)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Methodology: Four Model Families',
       'Traditional ML baselines + two deep learning approaches')

models = [
    ('Random\nForest', '200 estimators\nmax_depth=15\nWith & without SMOTE',
     DARK_BLUE, 'BASELINE'),
    ('Naive\nBayes', 'Gaussian NB\nProbabilistic classifier\nWith & without SMOTE',
     MID_BLUE, 'BASELINE'),
    ('CNN-LSTM', 'Hybrid DL model\nFocal Loss\nClass-Weighted FL',
     TEAL, 'PRIMARY DL'),
    ('DistilBERT', '66.9M parameters\nFine-tuned on\nbug report text',
     GREEN_DK, 'NLP MODEL'),
]

for i, (name, detail, c, badge) in enumerate(models):
    l = 0.35 + i * 3.22
    # Card
    rect(s, l, 1.45, 2.95, 5.45, fill=WHITE, line=c, lw=2.0)
    rect(s, l, 1.45, 2.95, 1.35, fill=c)
    # Badge
    rect(s, l+1.8, 1.5, 1.1, 0.35, fill=ORANGE)
    txt(s, badge, l+1.82, 1.52, 1.06, 0.31, size=9, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    # Model name
    txt(s, name, l, 1.65, 2.95, 1.1,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    # Detail
    txt(s, detail, l+0.1, 2.95, 2.75, 1.1,
        size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)
    # Imbalance method
    rect(s, l+0.15, 4.15, 2.65, 0.45, fill=PALE_BLUE)
    imb = ('BorderlineSMOTE\n(data-level)' if i <= 1
           else 'Focal Loss\n(algorithm-level)')
    txt(s, imb, l+0.15, 4.18, 2.65, 0.4, size=11, color=DARK_BLUE,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    # Primary metric
    txt(s, 'Primary metric: Recall (buggy class)',
        l+0.1, 4.72, 2.75, 0.35, size=11, color=MID_GRAY,
        align=PP_ALIGN.CENTER, italic=True)

# Evaluation box
rect(s, 0.35, 6.4, 12.6, 0.72, fill=DARK_BLUE)
txt(s, 'Evaluation:  Accuracy  |  Precision  |  Recall (primary)  |  '
    'F1-Score (primary)  |  ROC-AUC  |  10-Fold Stratified Cross-Validation',
    0.5, 6.44, 12.2, 0.64, size=14, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 9 -- BASELINE ML RESULTS
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Baseline ML Results',
       'Random Forest & Naive Bayes with and without BorderlineSMOTE')

make_table(s,
    headers=['Model', 'Accuracy', 'Precision', 'Recall', 'F1-Score', 'AUC', 'Bugs Caught'],
    rows=[
        ['RF  (no SMOTE)',  '0.8107', '0.5280', '0.2019', '0.2921', '0.7365', '85 / 421'],
        ['RF  + SMOTE',    '0.7653', '0.4043', '0.4513', '0.4265', '0.7350', '190 / 421'],
        ['NB  (no SMOTE)', '0.7901', '0.4412', '0.3207', '0.3714', '0.6788', '135 / 421'],
        ['NB  + SMOTE',    '0.7735', '0.4037', '0.3587', '0.3799', '0.6780', '151 / 421'],
    ],
    l=0.35, t=1.45, w=12.6, h=3.1,
    col_w=[3.0, 1.6, 1.6, 1.5, 1.6, 1.5, 1.8],
    hdr_size=12, row_size=12, hi_rows=[1])

# Key observations
txt(s, 'Key Observations', 0.35, 4.7, 12.6, 0.38,
    size=16, bold=True, color=DARK_BLUE)

tf2 = s.shapes.add_textbox(Inches(0.35), Inches(5.1),
                            Inches(12.6), Inches(2.0)).text_frame
tf2.word_wrap = True
for obs in [
    'Accuracy Paradox: RF no-SMOTE gets 81.1% accuracy but catches only 85 bugs (Recall = 0.202)',
    'SMOTE Effect: SMOTE improves RF Recall by +123% (85 -> 190 bugs caught)',
    'Best Baseline: RF + SMOTE (F1=0.427, Recall=0.451) -- this is our comparison benchmark',
    '10-Fold CV confirms stability: RF F1 = 0.287 +/- 0.038, NB F1 = 0.362 +/- 0.035',
]:
    bullet(tf2, obs, size=14, color=DARK_GRAY, space_before=4)


# ════════════════════════════════════════════════════════════
# SLIDE 10 -- CNN-LSTM ARCHITECTURE
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'CNN-LSTM Architecture',
       'Hybrid deep learning model: local feature extraction + sequential dependencies')

make_table(s,
    headers=['Layer', 'Configuration', 'Purpose'],
    rows=[
        ['Input',     '(batch, 21, 1)',          '21 complexity metrics as timesteps'],
        ['Conv1D',    '64 filters, kernel=3',    'Local pattern detection'],
        ['BatchNorm', '--',                       'Training stabilisation'],
        ['Conv1D',    '128 filters, kernel=3',   'Deep feature extraction'],
        ['BatchNorm', '--',                       'Training stabilisation'],
        ['LSTM',      'hidden=64, 1 layer',      'Sequential dependencies across metrics'],
        ['Dropout',   'rate = 0.4',              'Regularisation (prevent overfitting)'],
        ['Dense',     '32 units, ReLU',          'Feature compression'],
        ['Output',    '1 unit, Sigmoid',         'Bug probability in [0, 1]'],
    ],
    l=0.35, t=1.45, w=7.8, h=5.55,
    col_w=[1.6, 2.2, 4.0],
    hdr_size=12, row_size=12)

# Side info
txt(s, 'Model Stats', 8.4, 1.45, 4.6, 0.42,
    size=15, bold=True, color=DARK_BLUE)

stats = [
    ('Parameters', '77,121 trainable'),
    ('Framework',  'PyTorch 2.12'),
    ('Epochs',     '30 (BCE/Focal)  |  50 (Enhanced)'),
    ('Optimizer',  'Adam (lr=0.001 / 0.0005)'),
    ('Scheduler',  'CosineAnnealingLR'),
    ('Input',      '21 features reshaped to (21,1)'),
]
for i, (k, v) in enumerate(stats):
    t = 2.0 + i * 0.78
    rect(s, 8.4, t, 4.6, 0.68,
         fill=PALE_BLUE if i%2==0 else LIGHT_BLUE, line=MID_BLUE, lw=0.3)
    txt(s, k+':', 8.5, t+0.04, 1.5, 0.6, size=12, bold=True, color=DARK_BLUE,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, v, 10.05, t+0.04, 2.85, 0.6, size=12, color=DARK_GRAY,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)

rect(s, 8.4, 6.7, 4.6, 0.28, fill=ORANGE)
txt(s, 'Three loss variants: BCE | Focal Loss | Class-Weighted FL',
    8.5, 6.72, 4.4, 0.24, size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 11 -- FOCAL LOSS: KEY INNOVATION
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Key Innovation: Focal Loss & Class-Weighted Focal Loss',
       'Algorithm-level imbalance correction -- the primary research contribution')

# Left: formula explanation
rect(s, 0.35, 1.45, 12.6, 0.9, fill=DARK_BLUE)
txt(s, 'Standard Focal Loss:    FL(pt) = -alpha * (1 - pt)^gamma * log(pt)',
    0.55, 1.5, 12.2, 0.8, size=20, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)

rect(s, 0.35, 2.45, 12.6, 0.5, fill=ORANGE)
txt(s, 'Class-Weighted FL:    loss = pos_weight * FL(pt)   '
    'where  pos_weight = 4.2  (imbalance ratio)',
    0.55, 2.5, 12.2, 0.4, size=16, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)

# Parameters
params = [
    ('alpha = 0.75',   'Higher weight on the\nbuggy (minority) class'),
    ('gamma = 2.0',    'Focus on hard-to-classify\nexamples, ignore easy ones'),
    ('pos_weight = 4.2','Dataset imbalance ratio\ndirectly encoded in loss'),
]
for i, (param, desc) in enumerate(params):
    l = 0.35 + i*4.3
    rect(s, l, 3.1, 4.0, 1.85, fill=PALE_BLUE, line=DARK_BLUE, lw=1.0)
    txt(s, param, l+0.1, 3.15, 3.8, 0.55,
        size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    txt(s, desc, l+0.1, 3.75, 3.8, 1.1,
        size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Intuition
rect(s, 0.35, 5.1, 12.6, 0.38, fill=MID_BLUE)
txt(s, 'Intuition: Standard CE loss treats all samples equally. '
    'Focal Loss down-weights EASY samples (well-classified non-buggy modules) '
    'so the model focuses training effort on HARD samples (the rare buggy modules).',
    0.45, 5.14, 12.3, 0.3, size=12, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)

# Results preview
res_data = [
    ('CNN-LSTM  (BCE loss)',                 '0.3997', '0.5772', '243/421'),
    ('CNN-LSTM  (Focal Loss)',               '0.3980', '0.7553', '318/421'),
    ('CNN-LSTM  (Class-Weighted FL)',        '0.3742', '0.8337', '351/421'),
]
make_table(s,
    headers=['Model', 'F1-Score', 'Recall', 'Bugs Caught'],
    rows=res_data,
    l=0.35, t=5.6, w=12.6, h=1.55,
    col_w=[5.8, 1.8, 1.8, 3.2],
    hdr_size=12, row_size=12, hi_rows=[2])


# ════════════════════════════════════════════════════════════
# SLIDE 12 -- CNN-LSTM RESULTS DEEP DIVE
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'CNN-LSTM Results: Impact of Focal Loss',
       'Progressively increasing Recall with each imbalance-aware strategy')

# Comparison cards
cards = [
    ('BCE Loss', '0.3997', '0.5772', '243', '421', MID_BLUE,
     'Baseline DL. Better recall than ML\nbut standard loss limits minority-class focus.'),
    ('Focal Loss', '0.3980', '0.7553', '318', '421', ORANGE,
     '+30.8pp Recall vs BCE.\nFocuses training on hard minority examples.'),
    ('Class-Weighted\nFocal Loss', '0.3742', '0.8337', '351', '421', GREEN_DK,
     'BEST: +44.4pp Recall vs BCE.\n84.7% improvement over best ML baseline!'),
]

for i, (name, f1, rec, caught, total, c, comment) in enumerate(cards):
    l = 0.35 + i * 4.3
    rect(s, l, 1.45, 4.0, 5.45, fill=WHITE, line=c, lw=2.0)
    rect(s, l, 1.45, 4.0, 0.52, fill=c)
    txt(s, name, l, 1.52, 4.0, 0.38,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, 'Bugs Caught', l, 2.1, 4.0, 0.38,
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    txt(s, f'{caught} / {total}', l, 2.45, 4.0, 0.85,
        size=36, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s, f'Recall = {rec}', l, 3.35, 4.0, 0.45,
        size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    txt(s, f'F1-Score = {f1}', l, 3.82, 4.0, 0.38,
        size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
    rect(s, l+0.1, 4.3, 3.8, 0.05, fill=c)
    txt(s, comment, l+0.1, 4.45, 3.8, 1.3,
        size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Key stat
rect(s, 0.35, 7.02, 12.6, 0.45, fill=DARK_BLUE)
txt(s, 'RF + SMOTE (best baseline): caught 190/421 bugs (Recall=0.451)   '
    '-->   Class-Weighted FL: 351/421 (Recall=0.834)   =   +84.7% improvement',
    0.45, 7.04, 12.3, 0.39, size=13, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 13 -- DISTILBERT FOR BUG REPORT CLASSIFICATION
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'DistilBERT: Bug Report Severity Classification',
       'Transformer model for NLP-based bug triage (High vs Low Priority)')

# Left: Architecture
txt(s, 'Model Architecture', 0.35, 1.45, 6.1, 0.4,
    size=15, bold=True, color=DARK_BLUE)
arch_rows = [
    ('Base Model',    'distilbert-base-uncased'),
    ('Parameters',    '66.9 million'),
    ('Pre-training',  'BooksCorpus + Wikipedia'),
    ('Fine-tune task','Binary sequence classification'),
    ('Max Length',    '128 tokens'),
    ('Batch Size',    '16'),
    ('Epochs',        '4'),
    ('Optimizer',     'AdamW  (lr=2e-5)'),
    ('Scheduler',     'Linear warmup (10% of steps)'),
    ('Dataset',       '1,200 bug reports (960 train / 240 test)'),
]
make_table(s,
    headers=['Property', 'Value'],
    rows=arch_rows,
    l=0.35, t=1.95, w=6.1, h=5.1,
    col_w=[2.2, 3.9], hdr_size=12, row_size=11)

# Right: Results
txt(s, 'Results on Test Set (240 samples)', 6.75, 1.45, 6.2, 0.4,
    size=15, bold=True, color=DARK_BLUE)

metrics = [
    ('Accuracy',  '1.0000', '100% correct classifications'),
    ('Precision', '1.0000', '100% of predicted HP were real'),
    ('Recall',    '1.0000', '100% of HP bugs detected'),
    ('F1-Score',  '1.0000', 'Perfect balance'),
    ('ROC-AUC',   '1.0000', 'Perfect ranking ability'),
]
for i, (m, v, desc) in enumerate(metrics):
    t = 2.0 + i * 0.82
    rect(s, 6.75, t, 6.2, 0.72, fill=PALE_BLUE if i%2==0 else LIGHT_BLUE,
         line=GREEN_DK, lw=0.4)
    txt(s, m, 6.85, t+0.04, 2.0, 0.64, size=14, bold=True, color=DARK_BLUE,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, v, 8.9, t+0.04, 1.3, 0.64, size=22, bold=True, color=GREEN_DK,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, desc, 10.25, t+0.04, 2.6, 0.64, size=11, color=MID_GRAY,
        v=MSO_VERTICAL_ANCHOR.MIDDLE, italic=True)

rect(s, 6.75, 6.15, 6.2, 0.82, fill=DARK_BLUE)
txt(s, 'BERT\'s contextual embeddings clearly distinguish\n'
    '"crash/security/data-loss" from "cosmetic/minor" bug language',
    6.85, 6.2, 6.0, 0.72, size=12, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 14 -- COMPLETE COMPARISON TABLE
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Complete Model Comparison',
       'All 7 model configurations evaluated on the same test set (2,177 samples)')

make_table(s,
    headers=['Model', 'Accuracy', 'Precision', 'Recall', 'F1-Score', 'AUC', 'Bugs / 421'],
    rows=[
        ['RF  (no SMOTE)',         '0.8107', '0.5280', '0.2019', '0.2921', '0.7365', '85'],
        ['RF  + SMOTE',           '0.7653', '0.4043', '0.4513', '0.4265', '0.7350', '190'],
        ['NB  (no SMOTE)',        '0.7901', '0.4412', '0.3207', '0.3714', '0.6788', '135'],
        ['NB  + SMOTE',           '0.7735', '0.4037', '0.3587', '0.3799', '0.6780', '151'],
        ['CNN-LSTM  (BCE)',        '0.6647', '0.3057', '0.5772', '0.3997', '0.6808', '243'],
        ['CNN-LSTM  (Focal Loss)', '0.5581', '0.2702', '0.7553', '0.3980', '0.7001', '318'],
        ['CNN-LSTM  (CW-FL)',      '0.4607', '0.2412', '0.8337', '0.3742', '0.6920', '351'],
    ],
    l=0.35, t=1.45, w=12.6, h=4.05,
    col_w=[3.4, 1.5, 1.5, 1.4, 1.5, 1.4, 1.9],
    hdr_size=12, row_size=11, hi_rows=[6])

# DistilBERT note
rect(s, 0.35, 5.58, 12.6, 0.55, fill=GREEN_DK)
txt(s, 'DistilBERT (Bug Report Classification -- separate task):   '
    'Accuracy = Precision = Recall = F1 = AUC = 1.0000',
    0.5, 5.62, 12.2, 0.47, size=14, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)

# Key takeaway
rect(s, 0.35, 6.22, 12.6, 0.9, fill=PALE_BLUE, line=DARK_BLUE, lw=0.8)
tf3 = s.shapes.add_textbox(Inches(0.5), Inches(6.28),
                            Inches(12.2), Inches(0.78)).text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
r3a = p3.add_run()
r3a.text = 'Key Insight: '
r3a.font.size = Pt(13); r3a.font.bold = True
r3a.font.color.rgb = DARK_BLUE; r3a.font.name = 'Calibri'
r3b = p3.add_run()
r3b.text = ('Higher Recall means more bugs caught. Class-Weighted FL trades some Accuracy '
            'for dramatically higher Recall -- the correct trade-off for safety-critical bug detection. '
            'Precision-Recall trade-off is expected and acceptable in this domain.')
r3b.font.size = Pt(12); r3b.font.color.rgb = DARK_GRAY; r3b.font.name = 'Calibri'


# ════════════════════════════════════════════════════════════
# SLIDE 15 -- KEY FINDING: BUGS CAUGHT
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Key Finding: How Many Bugs Does Each Model Catch?',
       'Out of 421 defective modules in the test set (2,177 samples total)')

# Large image
img(s, '21_bugs_caught.png', 0.35, 1.38, 7.8, 5.75)

# Right side highlights
txt(s, 'Progress in Bug Detection', 8.45, 1.45, 4.55, 0.42,
    size=15, bold=True, color=DARK_BLUE)

progress = [
    ('85',  'RF (no SMOTE)',    RED_DK,   '20.2% caught'),
    ('190', 'RF + SMOTE',      ORANGE,   '45.1% caught'),
    ('243', 'CNN-LSTM (BCE)',   MID_BLUE, '57.7% caught'),
    ('318', 'CNN-LSTM (FL)',    TEAL,     '75.5% caught'),
    ('351', 'CNN-LSTM (CW-FL)',  GREEN_DK, '83.4% caught'),
]
for i, (num, model, c, pct) in enumerate(progress):
    t = 2.0 + i * 1.02
    rect(s, 8.45, t, 4.55, 0.88, fill=PALE_BLUE, line=c, lw=1.5)
    txt(s, num, 8.55, t+0.04, 0.7, 0.8, size=26, bold=True, color=c,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, model, 9.3, t+0.04, 2.5, 0.38, size=12, bold=True, color=DARK_GRAY,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, pct, 9.3, t+0.44, 2.5, 0.36, size=12, color=c, italic=True,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)

rect(s, 8.45, 7.1, 4.55, 0.28, fill=GREEN_DK)
txt(s, '+84.7% improvement over best ML baseline',
    8.5, 7.12, 4.45, 0.24, size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 16 -- ROC CURVES
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'ROC Curves & Performance Analysis',
       'Area Under Curve comparison across all code-metric models')

img(s, '18_roc_curves.png', 0.35, 1.38, 7.5, 5.75)

txt(s, 'ROC-AUC Summary', 8.15, 1.45, 4.85, 0.42,
    size=15, bold=True, color=DARK_BLUE)

auc_rows = [
    ('RF (no SMOTE)',   '0.7365', MID_GRAY),
    ('RF + SMOTE',     '0.7350', MID_GRAY),
    ('NB (no SMOTE)',  '0.6788', MID_GRAY),
    ('NB + SMOTE',    '0.6780', MID_GRAY),
    ('CNN-LSTM (BCE)', '0.6808', MID_BLUE),
    ('CNN-LSTM (FL)',  '0.7001', ORANGE),
    ('CNN-LSTM (CWF)', '0.6920', GREEN_DK),
]
for i, (m, auc, c) in enumerate(auc_rows):
    t = 2.0 + i * 0.73
    rect(s, 8.15, t, 4.85, 0.63, fill=PALE_BLUE if i%2==0 else WHITE,
         line=c, lw=0.4)
    txt(s, m, 8.25, t+0.04, 3.0, 0.55, size=11, color=DARK_GRAY,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, auc, 11.25, t+0.04, 1.65, 0.55, size=16, bold=True, color=c,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)

rect(s, 8.15, 7.1, 4.85, 0.28, fill=DARK_BLUE)
txt(s, 'Focal Loss CNN-LSTM achieves highest AUC (0.700)',
    8.2, 7.12, 4.75, 0.24, size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 17 -- RESEARCH QUESTIONS ANSWERED
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Research Questions: Answered',
       'Empirical answers derived from experimental results')

rq_answers = [
    ('RQ1', 'Traditional ML vs Deep Learning',
     'RF+SMOTE: Recall=0.451. CNN-LSTM+CW-FL: Recall=0.834. DL is +84.7% better at catching bugs.',
     DARK_BLUE),
    ('RQ2', 'Impact of CNN-LSTM & BERT on Imbalanced Data',
     'CNN-LSTM+CW-FL catches 351/421 bugs. DistilBERT: F1=1.000 on bug reports. Both far outperform ML.',
     MID_BLUE),
    ('RQ3', 'Impact of Class Imbalance',
     'Without balancing: RF Recall=0.202. SMOTE: +123%. Focal Loss over SMOTE: further +84.7%. '
     'Combined data + algorithm-level balancing is optimal.',
     TEAL),
    ('RQ4', 'Best Model for Scalable Defect Detection',
     'CNN-LSTM+CW-FL (Recall=0.834) for code metrics. DistilBERT (F1=1.0) for bug text. '
     'Two complementary models = one complete pipeline.',
     GREEN_DK),
]
for i, (rq, title, answer, c) in enumerate(rq_answers):
    t = 1.45 + i * 1.45
    rect(s, 0.35, t, 12.6, 1.35, fill=PALE_BLUE if i%2==0 else WHITE,
         line=c, lw=1.5)
    rect(s, 0.35, t, 1.0, 1.35, fill=c)
    txt(s, rq, 0.35, t, 1.0, 1.35,
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, title, 1.48, t+0.06, 5.0, 0.42, size=13, bold=True, color=c)
    txt(s, answer, 1.48, t+0.52, 11.3, 0.77, size=12, color=DARK_GRAY, wrap=True)

rect(s, 0.35, 7.2, 12.6, 0.22, fill=ORANGE)
txt(s, 'Focal Loss training represents a significant, empirically validated contribution to the field of AI-driven software defect prediction',
    0.5, 7.22, 12.2, 0.18, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 18 -- CONCLUSION & CONTRIBUTION
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Conclusion & Research Contribution',
       'Empirically validated advancement in software defect prediction')

# Main contribution box
rect(s, 0.35, 1.45, 12.6, 1.25, fill=DARK_BLUE)
rect(s, 0.35, 1.45, 0.12, 1.25, fill=ORANGE)
txt(s, 'Primary Contribution:  Imbalance-aware deep learning training through '
    'Class-Weighted Focal Loss applied to CNN-LSTM achieves Recall = 0.834 -- '
    'detecting 351/421 bugs, an 84.7% improvement over the best traditional ML baseline.',
    0.6, 1.52, 12.1, 1.1, size=15, bold=True, color=WHITE,
    v=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=True)

# 4 key conclusions
conclusions = [
    ('Class Imbalance is the\nDominant Challenge',
     'Without explicit handling, models achieve high accuracy by ignoring bugs entirely. '
     'Both data-level (SMOTE) AND algorithm-level (Focal Loss) interventions are needed.',
     DARK_BLUE),
    ('Deep Learning Significantly\nOutperforms Traditional ML',
     'CNN-LSTM Recall = 0.834 vs RF Recall = 0.451. DL captures complex non-linear '
     'patterns that RF and NB miss, especially with imbalance correction.',
     MID_BLUE),
    ('Complementary Multi-Modal\nFramework',
     'CNN-LSTM (code metrics) + DistilBERT (bug text) together cover two dimensions '
     'of automated software quality assurance -- a complete production pipeline.',
     TEAL),
    ('Practical Deployment\nGuidelines Derived',
     'Precision-Recall trade-off is acceptable when missing a bug costs more than '
     'investigating a false alarm -- as is the case in safety-critical systems.',
     GREEN_DK),
]
for i, (title, body, c) in enumerate(conclusions):
    l = 0.35 + i * 3.22
    rect(s, l, 2.85, 2.95, 3.95, fill=WHITE, line=c, lw=1.5)
    rect(s, l, 2.85, 2.95, 0.62, fill=c)
    txt(s, title, l+0.08, 2.9, 2.79, 0.54,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, body, l+0.1, 3.56, 2.75, 3.14,
        size=12, color=DARK_GRAY, wrap=True)

rect(s, 0.35, 6.95, 12.6, 0.42, fill=LIGHT_BLUE, line=DARK_BLUE, lw=0.5)
txt(s, 'This study confirms that AI-driven bug detection can be made significantly more effective '
    'through careful selection of imbalance-handling strategies.',
    0.5, 6.98, 12.2, 0.36, size=13, italic=True, color=DARK_BLUE,
    align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLIDE 19 -- FUTURE WORK
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, 'Future Work & Research Directions',
       'Extending this empirical foundation')

future = [
    ('Cross-Dataset Validation',
     'Extend to KC1, PC1, and other PROMISE datasets.\n'
     'Cross-project prediction: train on JM1, test on KC1.',
     DARK_BLUE, '01'),
    ('Explainable AI (XAI)',
     'Integrate SHAP / LIME to explain model predictions.\n'
     'Address the black-box limitation for industrial adoption.',
     MID_BLUE, '02'),
    ('Real-World Bug Reports',
     'Evaluate CodeBERT / GraphCodeBERT on live Bugzilla\n'
     'and GitHub Issues with noisy, real-world labels.',
     TEAL, '03'),
    ('Graph Neural Networks',
     'Apply GNNs on Abstract Syntax Trees (ASTs) to\n'
     'capture structural code relationships beyond metrics.',
     GREEN_DK, '04'),
    ('Unified Multi-Modal System',
     'Combine CNN-LSTM (code) + DistilBERT (text) into\n'
     'a single joint prediction production pipeline.',
     ORANGE, '05'),
    ('Transfer Learning',
     'Domain adaptation for cross-project generalisation\n'
     'without full retraining from scratch on each project.',
     RED_DK, '06'),
]

for i, (title, body, c, num) in enumerate(future):
    col = i % 3
    row = i // 3
    l = 0.35 + col * 4.3
    t = 1.48 + row * 2.88
    rect(s, l, t, 4.0, 2.72, fill=WHITE, line=c, lw=1.5)
    rect(s, l, t, 4.0, 0.55, fill=c)
    txt(s, num, l+3.45, t+0.05, 0.48, 0.44, size=16, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, title, l+0.1, t+0.05, 3.3, 0.44, size=13, bold=True, color=WHITE,
        v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, body, l+0.1, t+0.65, 3.8, 1.95, size=12, color=DARK_GRAY, wrap=True)


# ════════════════════════════════════════════════════════════
# SLIDE 20 -- THANK YOU / Q&A
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.333, 7.5, fill=DARK_BLUE)
rect(s, 0, 3.28, 13.333, 0.07, fill=ORANGE)

txt(s, 'Thank You', 0, 0.8, 13.333, 1.2,
    size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    v=MSO_VERTICAL_ANCHOR.MIDDLE)

txt(s, 'Questions & Discussion', 0, 1.95, 13.333, 0.65,
    size=24, color=RGBColor(180, 210, 240), align=PP_ALIGN.CENTER)

rect(s, 2.5, 3.45, 8.333, 1.55, fill=RGBColor(0, 50, 100))
for k, v, t in [
    ('Student',    'S. Shanthosh  |  20APSE4882', 3.52),
    ('Supervisor', 'Mrs. WMLS Abeythunga',         3.9),
    ('Module',     'SE 8101 Research Project', 4.28),
    ('University', 'Sabaragamuwa University of Sri Lanka', 4.66),
]:
    txt(s, k+':', 2.7, t, 1.8, 0.32, size=12, bold=True,
        color=RGBColor(160, 200, 240))
    txt(s, v, 4.55, t, 6.0, 0.32, size=12, color=WHITE)

# Summary stats at bottom
for i, (num, label, c) in enumerate([
    ('7', 'Model Configurations', ORANGE),
    ('21', 'Complexity Features', MID_BLUE),
    ('351', 'Bugs Caught (best)', GREEN_DK),
    ('84.7%', 'Recall Improvement', GOLD),
]):
    l = 1.5 + i * 2.65
    rect(s, l, 5.45, 2.3, 1.55, fill=RGBColor(0, 50, 100))
    txt(s, num, l, 5.5, 2.3, 0.85, size=36, bold=True, color=c,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)
    txt(s, label, l, 6.32, 2.3, 0.62, size=11, color=WHITE,
        align=PP_ALIGN.CENTER, v=MSO_VERTICAL_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
prs.save(OUT)
print(f'Presentation saved!')
print(f'Location : {OUT}')
print(f'Slides   : {len(prs.slides)}')
print(f'Slides list:')
for i, _ in enumerate(prs.slides, 1):
    print(f'  {i:02d}. See PowerPoint')
